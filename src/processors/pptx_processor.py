"""
PPTX text processor for extracting slide content
"""

from pptx import Presentation
from typing import List, Dict, Any
import io


def extract_text_from_pptx(file_content) -> Dict[str, Any]:
    """
    Extract text from uploaded PPTX file
    Returns structured data with slide content
    """
    try:
        # Create presentation from file content
        prs = Presentation(io.BytesIO(file_content))
        
        slides_data = []
        
        for slide_idx, slide in enumerate(prs.slides, 1):
            slide_content = {
                "slide_number": slide_idx,
                "title": "",
                "content": [],
                "notes": ""
            }
            
            # Extract title
            if slide.shapes.title and slide.shapes.title.text.strip():
                slide_content["title"] = slide.shapes.title.text.strip()
            
            # Extract other text content
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    # Skip title (already captured)
                    if shape == slide.shapes.title:
                        continue
                    
                    text_content = shape.text.strip()
                    if text_content:
                        slide_content["content"].append(text_content)
            
            # Extract slide notes
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    slide_content["notes"] = notes_text
            
            slides_data.append(slide_content)
        
        return {
            "success": True,
            "total_slides": len(slides_data),
            "slides": slides_data
        }
        
    except Exception as e:
        return {
            "success": False,
            "error": str(e),
            "total_slides": 0,
            "slides": []
        }