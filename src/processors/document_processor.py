"""
Document processor for extracting text from PPTX and PDF files.
Converts documents to structured JSON format with attributes (style, location, size, text).
Each slide/page is numbered, with nested data for elements.
"""

import json
import re
from typing import Dict, List, Any, Optional
from dataclasses import dataclass, asdict
import logging
import io

# Document processing libraries
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import PyPDF2
import pdfplumber

logger = logging.getLogger(__name__)


@dataclass
class TextElement:
    """Represents a text element with complete metadata"""
    text: str
    element_type: str  # title, bullet, body, note
    slide_number: int
    element_id: str
    style: Optional[Dict[str, Any]] = None
    location: Optional[Dict[str, Any]] = None
    size: Optional[Dict[str, Any]] = None
    font_info: Optional[Dict[str, Any]] = None


class DocumentProcessor:
    """Main class for processing presentation documents following exact specifications"""
    
    def __init__(self):
        self.supported_formats = ['.pptx', '.pdf']
    
    def process_document(self, file_content: bytes, filename: str) -> Dict[str, Any]:
        """
        Process document and return structured JSON with attributes
        (style, location, size, text). Each slide/page numbered with nested data.
        """
        file_extension = filename.lower().split('.')[-1]
        
        if file_extension == 'pptx':
            return self._process_pptx(file_content, filename)
        elif file_extension == 'pdf':
            return self._process_pdf(file_content, filename)
        else:
            raise ValueError(f"Unsupported file format: {file_extension}")
    
    def _process_pptx(self, file_content: bytes, filename: str) -> Dict[str, Any]:
        """Process PowerPoint file following specification requirements"""
        try:
            prs = Presentation(io.BytesIO(file_content))
            elements = []
            
            for slide_idx, slide in enumerate(prs.slides, 1):
                # Process slide title
                if slide.shapes.title:
                    title_element = self._create_pptx_element(
                        slide.shapes.title, 
                        slide_idx, 
                        f"slide_{slide_idx}_title",
                        "title"
                    )
                    if title_element:
                        elements.append(title_element)
                
                # Process all other shapes
                for shape_idx, shape in enumerate(slide.shapes):
                    if hasattr(shape, 'text') and shape.text.strip():
                        # Skip title (already processed)
                        if shape == slide.shapes.title:
                            continue
                        
                        element_type = self._classify_pptx_element(shape)
                        element = self._create_pptx_element(
                            shape,
                            slide_idx,
                            f"slide_{slide_idx}_shape_{shape_idx}",
                            element_type
                        )
                        if element:
                            elements.append(element)
                
                # Process slide notes
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
                    notes_element = TextElement(
                        text=slide.notes_slide.notes_text_frame.text.strip(),
                        element_type="note",
                        slide_number=slide_idx,
                        element_id=f"slide_{slide_idx}_notes"
                    )
                    elements.append(notes_element)
            
            return {
                "document_type": "pptx",
                "filename": filename,
                "total_slides": len(prs.slides),
                "processing_status": "success",
                "elements": [asdict(elem) for elem in elements],
                "metadata": {
                    "slide_dimensions": self._get_slide_dimensions(prs),
                    "total_elements": len(elements)
                }
            }
            
        except Exception as e:
            logger.error(f"Error processing PPTX: {str(e)}")
            return {
                "document_type": "pptx",
                "filename": filename,
                "total_slides": 0,
                "processing_status": "error",
                "error_message": str(e),
                "elements": []
            }
    
    def _process_pdf(self, file_content: bytes, filename: str) -> Dict[str, Any]:
        """Process PDF file following specification requirements"""
        try:
            elements = []
            
            with pdfplumber.open(io.BytesIO(file_content)) as pdf:
                for page_idx, page in enumerate(pdf.pages, 1):
                    # Extract text with positioning
                    chars = page.chars
                    
                    if chars:
                        # Group characters into text blocks
                        text_blocks = self._group_pdf_characters(chars)
                        
                        for block_idx, block in enumerate(text_blocks):
                            element_type = self._classify_pdf_element(block['text'])
                            
                            text_element = TextElement(
                                text=block['text'],
                                element_type=element_type,
                                slide_number=page_idx,
                                element_id=f"page_{page_idx}_block_{block_idx}",
                                location={
                                    "x": block['x0'],
                                    "y": block['y0'],
                                    "x1": block['x1'],
                                    "y1": block['y1']
                                },
                                size={
                                    "width": block['x1'] - block['x0'],
                                    "height": block['y1'] - block['y0']
                                },
                                font_info={
                                    "size": block.get('size', 0),
                                    "font": block.get('fontname', 'Unknown')
                                }
                            )
                            elements.append(text_element)
            
            return {
                "document_type": "pdf",
                "filename": filename,
                "total_slides": len(pdf.pages) if 'pdf' in locals() else 0,
                "processing_status": "success",
                "elements": [asdict(elem) for elem in elements],
                "metadata": {
                    "total_elements": len(elements)
                }
            }
            
        except Exception as e:
            logger.error(f"Error processing PDF: {str(e)}")
            return {
                "document_type": "pdf",
                "filename": filename,
                "total_slides": 0,
                "processing_status": "error",
                "error_message": str(e),
                "elements": []
            }
    
    def _create_pptx_element(self, shape, slide_idx: int, element_id: str, element_type: str) -> Optional[TextElement]:
        """Create TextElement from PPTX shape with complete metadata"""
        try:
            text = shape.text.strip()
            if not text:
                return None
            
            return TextElement(
                text=text,
                element_type=element_type,
                slide_number=slide_idx,
                element_id=element_id,
                style=self._extract_pptx_style(shape),
                location=self._extract_pptx_location(shape),
                size=self._extract_pptx_size(shape),
                font_info=self._extract_pptx_font(shape)
            )
        except Exception:
            return None
    
    def _extract_pptx_style(self, shape) -> Dict[str, Any]:
        """Extract style information from PPTX shape"""
        style_info = {"shape_type": str(shape.shape_type)}
        
        try:
            if hasattr(shape, 'text_frame') and shape.text_frame.paragraphs:
                first_run = shape.text_frame.paragraphs[0].runs[0] if shape.text_frame.paragraphs[0].runs else None
                if first_run:
                    style_info.update({
                        "bold": first_run.font.bold,
                        "italic": first_run.font.italic,
                        "underline": first_run.font.underline
                    })
        except Exception:
            pass
        
        return style_info
    
    def _extract_pptx_location(self, shape) -> Dict[str, Any]:
        """Extract location information from PPTX shape"""
        try:
            return {
                "left": shape.left,
                "top": shape.top
            }
        except Exception:
            return {}
    
    def _extract_pptx_size(self, shape) -> Dict[str, Any]:
        """Extract size information from PPTX shape"""
        try:
            return {
                "width": shape.width,
                "height": shape.height
            }
        except Exception:
            return {}
    
    def _extract_pptx_font(self, shape) -> Dict[str, Any]:
        """Extract font information from PPTX shape"""
        font_info = {}
        try:
            if hasattr(shape, 'text_frame') and shape.text_frame.paragraphs:
                first_run = shape.text_frame.paragraphs[0].runs[0] if shape.text_frame.paragraphs[0].runs else None
                if first_run:
                    font_info.update({
                        "name": first_run.font.name,
                        "size": first_run.font.size.pt if first_run.font.size else None
                    })
        except Exception:
            pass
        
        return font_info
    
    def _classify_pptx_element(self, shape) -> str:
        """Classify PPTX element type"""
        try:
            text = shape.text.strip()
            
            # Check for bullet points
            if self._has_bullet_points(text):
                return "bullet"
            
            # Check if it's in a placeholder
            if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                return "body"
            
            # Default classification
            return "body"
            
        except Exception:
            return "body"
    
    def _classify_pdf_element(self, text: str) -> str:
        """Classify PDF text element by type"""
        text = text.strip()
        
        # Title detection heuristics
        if len(text) < 100 and (text.isupper() or text.istitle()):
            return "title"
        
        # Bullet detection
        if self._has_bullet_points(text):
            return "bullet"
        
        # Short text without period might be heading
        if len(text) < 80 and not text.endswith('.'):
            return "title"
        
        return "body"
    
    def _has_bullet_points(self, text: str) -> bool:
        """Detect if text contains bullet points"""
        bullet_patterns = [
            r'^\s*[•▪▫◦‣⁃]\s+',  # Unicode bullets
            r'^\s*[-*+]\s+',      # ASCII bullets
            r'^\s*\d+\.\s+',      # Numbered lists
            r'^\s*[a-zA-Z]\.\s+'  # Lettered lists
        ]
        
        for pattern in bullet_patterns:
            if re.search(pattern, text, re.MULTILINE):
                return True
        
        return False
    
    def _group_pdf_characters(self, chars: List[Dict]) -> List[Dict]:
        """Group PDF characters into text blocks"""
        if not chars:
            return []
        
        # Simple grouping by proximity (can be enhanced)
        blocks = []
        current_block = {
            'text': '',
            'x0': float('inf'),
            'y0': float('inf'),
            'x1': -float('inf'),
            'y1': -float('inf'),
            'size': 0,
            'fontname': ''
        }
        
        for char in chars:
            current_block['text'] += char['text']
            current_block['x0'] = min(current_block['x0'], char['x0'])
            current_block['y0'] = min(current_block['y0'], char['y0'])
            current_block['x1'] = max(current_block['x1'], char['x1'])
            current_block['y1'] = max(current_block['y1'], char['y1'])
            current_block['size'] = max(current_block['size'], char['size'])
            if not current_block['fontname']:
                current_block['fontname'] = char.get('fontname', 'Unknown')
        
        if current_block['text'].strip():
            blocks.append(current_block)
        
        return blocks
    
    def _get_slide_dimensions(self, presentation) -> Dict[str, int]:
        """Get slide dimensions from presentation"""
        try:
            return {
                "width": presentation.slide_width,
                "height": presentation.slide_height
            }
        except Exception:
            return {"width": 0, "height": 0}


# Helper function for easy import
def process_document(file_content: bytes, filename: str) -> Dict[str, Any]:
    """Convenience function to process a document"""
    processor = DocumentProcessor()
    return processor.process_document(file_content, filename)