"""
Document processor for extracting text from PPTX and PDF files.
Converts documents to structured JSON format with attributes (style, location, size, text).
Each slide/page is numbered, with nested data for elements.
"""

import json
import re
from typing import Dict, List, Any, Optional
from dataclasses import dataclass, asdict
import logging
import io

# Document processing libraries
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import PyPDF2
import pdfplumber

logger = logging.getLogger(__name__)


@dataclass
class TextElement:
    """Represents a text element with complete metadata"""
    text: str
    element_type: str  # title, bullet, body, note, author, caption
    slide_number: int
    element_id: str
    style: Optional[Dict[str, Any]] = None
    location: Optional[Dict[str, Any]] = None
    size: Optional[Dict[str, Any]] = None
    font_info: Optional[Dict[str, Any]] = None


class DocumentProcessor:
    """Main class for processing presentation documents following exact specifications"""
    
    def __init__(self):
        self.supported_formats = ['.pptx', '.pdf']
    
    def process_document(self, file_content: bytes, filename: str) -> Dict[str, Any]:
        """
        Process document and return structured JSON with attributes
        (style, location, size, text). Each slide/page numbered with nested data.
        """
        file_extension = filename.lower().split('.')[-1]
        
        if file_extension == 'pptx':
            return self._process_pptx(file_content, filename)
        elif file_extension == 'pdf':
            return self._process_pdf(file_content, filename)
        else:
            raise ValueError(f"Unsupported file format: {file_extension}")
    
    def _process_pptx(self, file_content: bytes, filename: str) -> Dict[str, Any]:
        """Process PowerPoint file following specification requirements"""
        try:
            prs = Presentation(io.BytesIO(file_content))
            elements = []
            
            for slide_idx, slide in enumerate(prs.slides, 1):
                # Process slide title
                if slide.shapes.title:
                    title_element = self._create_pptx_element(
                        slide.shapes.title, 
                        slide_idx, 
                        f"slide_{slide_idx}_title",
                        "title"
                    )
                    if title_element:
                        elements.append(title_element)
                
                # Process all other shapes
                for shape_idx, shape in enumerate(slide.shapes):
                    if hasattr(shape, 'text') and shape.text.strip():
                        # Skip title (already processed)
                        if shape == slide.shapes.title:
                            continue
                        
                        element_type = self._classify_pptx_element(shape)
                        element = self._create_pptx_element(
                            shape,
                            slide_idx,
                            f"slide_{slide_idx}_shape_{shape_idx}",
                            element_type
                        )
                        if element:
                            elements.append(element)
                
                # Process slide notes
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
                    notes_element = TextElement(
                        text=slide.notes_slide.notes_text_frame.text.strip(),
                        element_type="note",
                        slide_number=slide_idx,
                        element_id=f"slide_{slide_idx}_notes"
                    )
                    elements.append(notes_element)
            
            return {
                "document_type": "pptx",
                "filename": filename,
                "total_slides": len(prs.slides),
                "processing_status": "success",
                "elements": [asdict(elem) for elem in elements],
                "metadata": {
                    "slide_dimensions": self._get_slide_dimensions(prs),
                    "total_elements": len(elements)
                }
            }
            
        except Exception as e:
            logger.error(f"Error processing PPTX: {str(e)}")
            return {
                "document_type": "pptx",
                "filename": filename,
                "total_slides": 0,
                "processing_status": "error",
                "error_message": str(e),
                "elements": []
            }
    
    def _process_pdf(self, file_content: bytes, filename: str) -> Dict[str, Any]:
        """Process PDF with layout-aware extraction (bbox + font info)."""
        try:
            elements = []

            with pdfplumber.open(io.BytesIO(file_content)) as pdf:
                for page_idx, page in enumerate(pdf.pages, 1):
                    # words with coordinates (points; 1pt = 1/72 inch)
                    words = page.extract_words(
                        use_text_flow=True, keep_blank_chars=False,
                        x_tolerance=2, y_tolerance=3
                    )
                    if not words:
                        continue

                    # group words into lines by ~same 'top'
                    from collections import defaultdict, Counter
                    lines = defaultdict(list)
                    for w in words:
                        lines[round(w["top"], 1)].append(w)

                    # iterate lines
                    for order, (top, ws) in enumerate(sorted(lines.items(), key=lambda kv: kv[0])):
                        ws_sorted = sorted(ws, key=lambda w: w["x0"])
                        text = " ".join(w["text"] for w in ws_sorted).strip()

                        # filter junk like bare "undefined"
                        if not text or text.lower() == "undefined":
                            continue

                        x0 = min(w["x0"] for w in ws_sorted)
                        x1 = max(w["x1"] for w in ws_sorted)
                        y0 = min(w["top"] for w in ws_sorted)
                        y1 = max(w["bottom"] for w in ws_sorted)

                        # collect fonts/sizes from chars within bbox (tolerant mask)
                        tol = 0.5
                        chars = [
                            c for c in page.chars
                            if (c["x0"] >= x0 - tol and c["x1"] <= x1 + tol and
                                c["top"] >= y0 - tol and c["bottom"] <= y1 + tol)
                        ]
                        sizes = [c.get("size") for c in chars if c.get("size") is not None]
                        fonts = [c.get("fontname") for c in chars if c.get("fontname")]

                        # heuristic style flags
                        is_bold = any("Bold" in (f or "") for f in fonts)
                        is_italic = any(("Italic" in (f or "")) or ("Oblique" in (f or "")) for f in fonts)

                        element_type = self._classify_pdf_element(text)

                        elements.append(TextElement(
                            text=text,
                            element_type=element_type,
                            slide_number=page_idx,
                            element_id=f"page_{page_idx}_line_{order}",
                            location={"x": round(x0, 2), "y": round(y0, 2), "x1": round(x1, 2), "y1": round(y1, 2)},
                            size={"width": round(x1 - x0, 2), "height": round(y1 - y0, 2)},
                            font_info={
                                "font_most_common": Counter(fonts).most_common(1)[0][0] if fonts else None,
                                "size_median": float(sorted(sizes)[len(sizes)//2]) if sizes else None,
                                "units": "pt"
                            },
                            style={"bold": is_bold, "italic": is_italic}
                        ))

            return {
                "document_type": "pdf",
                "filename": filename,
                "total_slides": len(pdf.pages) if 'pdf' in locals() else 0,
                "processing_status": "success",
                "elements": [asdict(e) for e in elements],
                "metadata": {"total_elements": len(elements)}
            }

        except Exception as e:
            logger.error(f"Error processing PDF: {str(e)}")
            return {
                "document_type": "pdf",
                "filename": filename,
                "total_slides": 0,
                "processing_status": "error",
                "error_message": str(e),
                "elements": []
            }

        
    def _split_pdf_text_elements(self, text: str) -> List[str]:
        """Split PDF text into individual logical elements"""
        elements = []
        
        # Split by multiple newlines first (major sections)
        major_sections = re.split(r'\n\s*\n\s*', text.strip())
        
        for section in major_sections:
            section = section.strip()
            if not section:
                continue
            
            # Check if section contains bullet points
            if self._has_bullet_points(section):
                # Split bullet lists into individual bullets
                bullet_items = self._extract_bullet_items(section)
                elements.extend(bullet_items)
            else:
                # Check if section contains multiple sentences that should be split
                lines = [line.strip() for line in section.split('\n') if line.strip()]
                
                for line in lines:
                    # Each significant line becomes its own element
                    if len(line) > 10:  # Avoid very short fragments
                        elements.append(line)
        
        return elements
    
    def _extract_bullet_items(self, text: str) -> List[str]:
        """Extract individual bullet items from text"""
        bullets = []
        
        # Pattern for bullet points
        bullet_patterns = [
            r'^\s*[•▪▫◦‣⁃]\s*(.+)',     # Unicode bullets
            r'^\s*[-*+]\s*(.+)',         # ASCII bullets  
            r'^\s*\d+\.\s*(.+)',         # Numbered lists
            r'^\s*[a-zA-Z]\.\s*(.+)'     # Lettered lists
        ]
        
        lines = text.split('\n')
        current_bullet = ""
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
                
            # Check if line starts with a bullet
            is_bullet_start = False
            for pattern in bullet_patterns:
                if re.match(pattern, line):
                    is_bullet_start = True
                    break
            
            if is_bullet_start:
                # Save previous bullet if exists
                if current_bullet.strip():
                    bullets.append(current_bullet.strip())
                current_bullet = line
            else:
                # Continue previous bullet
                if current_bullet:
                    current_bullet += " " + line
                else:
                    # Standalone line
                    bullets.append(line)
        
        # Don't forget the last bullet
        if current_bullet.strip():
            bullets.append(current_bullet.strip())
        
        return bullets
    
    def _create_pptx_element(self, shape, slide_idx: int, element_id: str, element_type: str) -> Optional[TextElement]:
        """Create TextElement from PPTX shape with complete metadata"""
        try:
            text = shape.text.strip()
            if not text:
                return None
            
            return TextElement(
                text=text,
                element_type=element_type,
                slide_number=slide_idx,
                element_id=element_id,
                style=self._extract_pptx_style(shape),
                location=self._extract_pptx_location(shape),
                size=self._extract_pptx_size(shape),
                font_info=self._extract_pptx_font(shape)
            )
        except Exception:
            return None
    
    def _extract_pptx_style(self, shape) -> Dict[str, Any]:
        """Extract style information from PPTX shape"""
        style_info = {"shape_type": str(shape.shape_type)}
        
        try:
            if hasattr(shape, 'text_frame') and shape.text_frame.paragraphs:
                first_run = shape.text_frame.paragraphs[0].runs[0] if shape.text_frame.paragraphs[0].runs else None
                if first_run:
                    style_info.update({
                        "bold": first_run.font.bold,
                        "italic": first_run.font.italic,
                        "underline": first_run.font.underline
                    })
        except Exception:
            pass
        
        return style_info
    
    def _extract_pptx_location(self, shape) -> Dict[str, Any]:
        """Extract location information from PPTX shape"""
        try:
            return {
                "left": shape.left,
                "top": shape.top
            }
        except Exception:
            return {}
    
    def _extract_pptx_size(self, shape) -> Dict[str, Any]:
        """Extract size information from PPTX shape"""
        try:
            return {
                "width": shape.width,
                "height": shape.height
            }
        except Exception:
            return {}
    
    def _extract_pptx_font(self, shape) -> Dict[str, Any]:
        """Extract font information from PPTX shape"""
        font_info = {}
        try:
            if hasattr(shape, 'text_frame') and shape.text_frame.paragraphs:
                first_run = shape.text_frame.paragraphs[0].runs[0] if shape.text_frame.paragraphs[0].runs else None
                if first_run:
                    font_info.update({
                        "name": first_run.font.name,
                        "size": first_run.font.size.pt if first_run.font.size else None
                    })
        except Exception:
            pass
        
        return font_info
    
    def _classify_pptx_element(self, shape) -> str:
        """Classify PPTX element type"""
        try:
            text = shape.text.strip()
            
            # Check for bullet points
            if self._has_bullet_points(text):
                return "bullet"
            
            # Check if it's in a placeholder
            if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                return "body"
            
            # Default classification
            return "body"
            
        except Exception:
            return "body"
    
    def _classify_pdf_element(self, text: str) -> str:
        """Classify PDF text element by type with improved logic"""
        text = text.strip()
        
        # Title detection - improved heuristics
        if self._is_likely_title(text):
            return "title"
        
        # Bullet detection
        if self._has_bullet_points(text):
            return "bullet"
        
        # Author/affiliation detection
        if self._is_author_or_affiliation(text):
            return "author"
        
        # Caption or short descriptor
        if len(text) < 150 and (
            text.startswith(('Figure', 'Chart', 'Graph', 'Image', 'Photo', 'Source:')) or
            re.search(r'\d{4}|\%|°[CF]', text)  # Contains years, percentages, or temperatures
        ):
            return "caption"
        
        # Default to body text
        return "body"
    
    def _is_likely_title(self, text: str) -> bool:
        """Improved title detection"""
        text = text.strip()
        
        # Very short text (likely fragments)
        if len(text) < 5:
            return False
        
        # Check various title indicators
        title_indicators = [
            len(text) < 100 and text.isupper(),  # Short uppercase text
            len(text) < 80 and text.istitle() and not text.endswith('.'),  # Title case without period
            len(text) < 60 and text.count(' ') <= 8,  # Short with few words
            re.match(r'^[A-Z][^.!?]*$', text) and len(text) < 100  # Starts with capital, no sentence punctuation
        ]
        
        return any(title_indicators)
    
    def _is_author_or_affiliation(self, text: str) -> bool:
        """Detect author names or institutional affiliations"""
        text = text.strip()
        
        # Common patterns for authors and affiliations
        author_patterns = [
            r'^[A-Z][a-z]+ [A-Z][a-z]+$',  # First Last
            r'^[A-Z]\. [A-Z][a-z]+$',      # F. Last
            r'University|Institute|College|Department',  # Institutions
            r'@[a-zA-Z]+\.[a-zA-Z]+',      # Email
        ]
        
        return any(re.search(pattern, text) for pattern in author_patterns)
    
    def _has_bullet_points(self, text: str) -> bool:
        """Detect if text contains bullet points"""
        bullet_patterns = [
            r'^\s*[•▪▫◦‣⁃]\s+',  # Unicode bullets
            r'^\s*[-*+]\s+',      # ASCII bullets
            r'^\s*\d+\.\s+',      # Numbered lists
            r'^\s*[a-zA-Z]\.\s+'  # Lettered lists
        ]
        
        for pattern in bullet_patterns:
            if re.search(pattern, text, re.MULTILINE):
                return True
        
        return False
    
    def _get_slide_dimensions(self, presentation) -> Dict[str, int]:
        """Get slide dimensions from presentation"""
        try:
            return {
                "width": presentation.slide_width,
                "height": presentation.slide_height
            }
        except Exception:
            return {"width": 0, "height": 0}


# Helper function for easy import
def process_document(file_content: bytes, filename: str) -> Dict[str, Any]:
    """Convenience function to process a document"""
    processor = DocumentProcessor()
    return processor.process_document(file_content, filename)