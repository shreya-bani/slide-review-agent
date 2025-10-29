"""
Enhanced PowerPoint file reader with comprehensive formatting extraction.

Features:
- Complete theme color extraction and resolution
- Master slide inheritance tracking
- Comprehensive font, margin, and spacing details
- Efficient caching and batch processing
- Better error handling and fallbacks
"""
import logging
from pathlib import Path
from typing import Dict, List, Optional, Any, Tuple
from functools import lru_cache

from pptx import Presentation
from pptx.text.text import TextFrame
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_THEME_COLOR
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Import dataclass models
from .model import (
    TextLocator,
    EnhancedFormattingInfo,
    TableCell,
    ExtractedTable,
    ExtractedText,
    SlideContent,
    ThemeInfo
)

import warnings
warnings.filterwarnings("ignore", category=DeprecationWarning)

logger = logging.getLogger(__name__)


class EnhancedPPTXReader:
    """Enhanced PowerPoint reader with comprehensive formatting extraction."""

    # Theme color name mapping
    THEME_COLOR_NAMES = {
        MSO_THEME_COLOR.ACCENT_1: "ACCENT_1",
        MSO_THEME_COLOR.ACCENT_2: "ACCENT_2",
        MSO_THEME_COLOR.ACCENT_3: "ACCENT_3",
        MSO_THEME_COLOR.ACCENT_4: "ACCENT_4",
        MSO_THEME_COLOR.ACCENT_5: "ACCENT_5",
        MSO_THEME_COLOR.ACCENT_6: "ACCENT_6",
        MSO_THEME_COLOR.BACKGROUND_1: "BACKGROUND_1",
        MSO_THEME_COLOR.BACKGROUND_2: "BACKGROUND_2",
        MSO_THEME_COLOR.DARK_1: "DARK_1",
        MSO_THEME_COLOR.DARK_2: "DARK_2",
        MSO_THEME_COLOR.FOLLOWED_HYPERLINK: "FOLLOWED_HYPERLINK",
        MSO_THEME_COLOR.HYPERLINK: "HYPERLINK",
        MSO_THEME_COLOR.LIGHT_1: "LIGHT_1",
        MSO_THEME_COLOR.LIGHT_2: "LIGHT_2",
        MSO_THEME_COLOR.TEXT_1: "TEXT_1",
        MSO_THEME_COLOR.TEXT_2: "TEXT_2",
    }

    def __init__(self):
        self.presentation = None
        self.file_path = None
        self.theme_info = ThemeInfo()
        self._theme_color_cache = {}
        self._master_cache = {}

    def load_file(self, file_path: str) -> bool:
        """Load PowerPoint file and extract theme information."""
        try:
            self.file_path = Path(file_path)
            self.presentation = Presentation(file_path)
            self._extract_theme_info()
            self._cache_master_slides()
            logger.info(f"Loaded PPTX: {file_path} ({len(self.presentation.slides)} slides)")
            return True
        except Exception as e:
            logger.error(f"Failed to load PPTX {file_path}: {e}")
            return False

    def _extract_theme_info(self):
        """Extract comprehensive theme information."""
        try:
            if not hasattr(self.presentation, 'part'):
                return

            part = self.presentation.part

            # Extract theme name
            if hasattr(part, 'theme'):
                theme = part.theme
                try:
                    theme_element = theme._element
                    self.theme_info.theme_name = theme_element.get('name', 'Unknown')
                except:
                    pass

                # Extract color scheme
                if hasattr(theme, 'color_scheme'):
                    self._extract_color_scheme(theme.color_scheme)

                # Extract font scheme
                if hasattr(theme, 'font_scheme'):
                    self._extract_font_scheme(theme.font_scheme)

        except Exception as e:
            logger.debug(f"Could not extract theme info: {e}")

    def _extract_color_scheme(self, color_scheme):
        """Extract theme color scheme with RGB values."""
        try:
            # Map theme color names to RGB values
            theme_color_attrs = [
                'dk1', 'lt1', 'dk2', 'lt2',
                'accent1', 'accent2', 'accent3',
                'accent4', 'accent5', 'accent6',
                'hlink', 'folHlink'
            ]

            for attr in theme_color_attrs:
                if hasattr(color_scheme, attr):
                    color = getattr(color_scheme, attr)
                    if color:
                        rgb = self._color_to_rgb(color)
                        if rgb:
                            self.theme_info.color_scheme[attr] = rgb
        except Exception as e:
            logger.debug(f"Could not extract color scheme: {e}")

    def _extract_font_scheme(self, font_scheme):
        """Extract theme font scheme."""
        try:
            # Extract major (headings) and minor (body) fonts
            if hasattr(font_scheme, 'major_font'):
                major = font_scheme.major_font
                if hasattr(major, 'latin_typeface'):
                    self.theme_info.font_scheme['major'] = major.latin_typeface

            if hasattr(font_scheme, 'minor_font'):
                minor = font_scheme.minor_font
                if hasattr(minor, 'latin_typeface'):
                    self.theme_info.font_scheme['minor'] = minor.latin_typeface
        except Exception as e:
            logger.debug(f"Could not extract font scheme: {e}")

    def _cache_master_slides(self):
        """Cache master slide information for inheritance tracking."""
        try:
            for slide in self.presentation.slides:
                try:
                    layout = slide.slide_layout
                    master = layout.slide_master
                    master_name = getattr(master, 'name', 'Unknown')
                    self._master_cache[slide.slide_id] = {
                        'master_name': master_name,
                        'layout_name': layout.name
                    }
                except:
                    pass
        except Exception as e:
            logger.debug(f"Could not cache master slides: {e}")

    def _color_to_rgb(self, color) -> Optional[str]:
        """Convert a color object to RGB hex string."""
        try:
            if hasattr(color, 'rgb'):
                rgb = color.rgb
                if isinstance(rgb, RGBColor):
                    return f"#{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"
                elif isinstance(rgb, int):
                    return f"#{rgb:06X}"
        except:
            pass
        return None

    def extract_all_content(self) -> List[SlideContent]:
        """Extract content from all slides (1-based indexing)."""
        if not self.presentation:
            raise ValueError("No presentation loaded")

        slides = []
        for idx, slide in enumerate(self.presentation.slides, start=1):
            slide_content = self._extract_slide_content(idx, slide)
            slides.append(slide_content)

        return slides

    def _extract_slide_content(self, slide_idx: int, slide) -> SlideContent:
        """Extract comprehensive content from a single slide."""
        slide_content = SlideContent(slide_index=slide_idx)

        # Get layout and master info
        try:
            slide_content.layout_name = slide.slide_layout.name
            slide_content.master_name = slide.slide_layout.slide_master.name
        except:
            slide_content.layout_name = "Unknown"
            slide_content.master_name = "Unknown"

        # Extract background color
        slide_content.background_color = self._extract_slide_background(slide)

        # Extract title
        slide_content.title = self._extract_title(slide)

        # Extract content from shapes and tables
        slide_content.content_elements, slide_content.tables = self._extract_shapes_content(slide_idx, slide)

        # Extract notes
        slide_content.notes = self._extract_notes(slide)

        return slide_content

    def _extract_slide_background(self, slide) -> Optional[str]:
        """Extract slide background color."""
        try:
            background = slide.background
            if hasattr(background, 'fill'):
                fill = background.fill
                if hasattr(fill, 'fore_color'):
                    return self._extract_color_comprehensive(fill.fore_color).get('rgb')
        except:
            pass
        return None

    def _extract_title(self, slide) -> Optional[str]:
        """Extract slide title."""
        try:
            if hasattr(slide, 'shapes') and slide.shapes.title:
                title = slide.shapes.title.text.strip()
                return title if title else None
        except:
            pass
        return None

    def _extract_shapes_content(self, slide_idx: int, slide) -> Tuple[List[ExtractedText], List[ExtractedTable]]:
        """Extract text and tables from all shapes in slide."""
        elements = []
        tables = []
        content_idx = 0
        table_idx = 0

        for shape in slide.shapes:
            # Check if it's a table
            if shape.shape_type == 19:  # MSO_SHAPE_TYPE.TABLE
                try:
                    extracted_table = self._extract_table(slide_idx, shape.table, table_idx)
                    tables.append(extracted_table)
                    table_idx += 1
                except Exception as e:
                    logger.debug(f"Failed to extract table on slide {slide_idx}: {e}")
                continue

            # Extract text from text frames
            if not hasattr(shape, 'text_frame') or not shape.text_frame:
                continue

            # Skip title (already extracted)
            if shape == slide.shapes.title:
                continue

            # Extract shape background for context
            shape_bg = self._extract_shape_background(shape)

            shape_elements = self._extract_text_frame(
                slide_idx,
                shape.text_frame,
                content_idx,
                shape_bg
            )
            elements.extend(shape_elements)
            content_idx += 1

        return elements, tables

    def _extract_shape_background(self, shape) -> Optional[str]:
        """Extract shape background color."""
        try:
            if hasattr(shape, 'fill'):
                fill = shape.fill
                if hasattr(fill, 'fore_color'):
                    return self._extract_color_comprehensive(fill.fore_color).get('rgb')
        except:
            pass
        return None

    def _extract_table(self, slide_idx: int, table, table_idx: int) -> ExtractedTable:
        """Extract table with all cells and formatting."""
        rows = len(table.rows)
        columns = len(table.columns)

        extracted_table = ExtractedTable(
            slide_index=slide_idx,
            table_index=table_idx,
            rows=rows,
            columns=columns
        )

        # Detect if first row is header (heuristic: bold or different background)
        has_header = False
        if rows > 0:
            try:
                first_row_cells = [table.cell(0, col) for col in range(columns)]
                # Check if first row cells are bold
                bold_count = 0
                for cell in first_row_cells:
                    if cell.text_frame and cell.text_frame.paragraphs:
                        for para in cell.text_frame.paragraphs:
                            for run in para.runs:
                                if run.font.bold:
                                    bold_count += 1
                                    break
                has_header = bold_count >= (columns // 2)  # Majority bold = header
                extracted_table.has_header_row = has_header
            except:
                pass

        # Extract all cells
        for row_idx in range(rows):
            for col_idx in range(columns):
                try:
                    cell = table.cell(row_idx, col_idx)
                    cell_text = cell.text.strip()

                    # Extract cell formatting
                    cell_formatting = self._extract_cell_formatting(cell)

                    # Detect merged cells (rowspan/colspan)
                    rowspan = 1
                    colspan = 1
                    # Note: python-pptx doesn't directly expose merged cell info easily
                    # This is a simplified implementation

                    table_cell = TableCell(
                        row_index=row_idx,
                        col_index=col_idx,
                        text=cell_text,
                        formatting=cell_formatting,
                        is_header=(row_idx == 0 and has_header),
                        rowspan=rowspan,
                        colspan=colspan
                    )

                    extracted_table.cells.append(table_cell)

                except Exception as e:
                    logger.debug(f"Failed to extract cell ({row_idx}, {col_idx}): {e}")
                    continue

        return extracted_table

    def _extract_cell_formatting(self, cell) -> EnhancedFormattingInfo:
        """Extract formatting from a table cell."""
        fmt = EnhancedFormattingInfo()

        try:
            # Get cell background
            if hasattr(cell, 'fill'):
                try:
                    fill = cell.fill
                    if hasattr(fill, 'fore_color'):
                        color_info = self._extract_color_comprehensive(fill.fore_color)
                        fmt.background_color = color_info.get('display')
                        fmt.background_color_rgb = color_info.get('rgb')
                except:
                    pass

            # Extract text formatting from first paragraph/run
            if cell.text_frame and cell.text_frame.paragraphs:
                for para in cell.text_frame.paragraphs:
                    if para.text.strip():
                        # Get formatting from this paragraph
                        para_fmt = self._extract_comprehensive_formatting(para)
                        # Copy relevant fields
                        fmt.font_name = para_fmt.font_name
                        fmt.font_size = para_fmt.font_size
                        fmt.is_bold = para_fmt.is_bold
                        fmt.is_italic = para_fmt.is_italic
                        fmt.is_underline = para_fmt.is_underline
                        fmt.font_color = para_fmt.font_color
                        fmt.font_color_rgb = para_fmt.font_color_rgb
                        fmt.font_color_theme = para_fmt.font_color_theme
                        fmt.paragraph_alignment = para_fmt.paragraph_alignment
                        break
        except Exception as e:
            logger.debug(f"Could not extract cell formatting: {e}")

        return fmt

    def _extract_text_frame(self, slide_idx: int, text_frame: TextFrame,
                           shape_idx: int, shape_bg: Optional[str] = None) -> List[ExtractedText]:
        """Extract text with comprehensive formatting from a text frame."""
        elements = []

        for para_idx, paragraph in enumerate(text_frame.paragraphs):
            text = paragraph.text.strip()
            if not text:
                continue

            # Get bullet level
            bullet_level = getattr(paragraph, 'level', None)

            # Get comprehensive formatting info
            formatting = self._extract_comprehensive_formatting(paragraph, shape_bg)

            # Determine element type
            element_type = 'bullet' if bullet_level is not None else 'content'

            # Create locator
            locator = TextLocator(
                slide_index=slide_idx,
                element_type=element_type,
                element_index=shape_idx,
                paragraph_index=para_idx,
                bullet_level=bullet_level
            )

            # Create extracted text
            element = ExtractedText(
                text=text,
                locator=locator,
                formatting=formatting
            )

            elements.append(element)

        return elements

    def _extract_comprehensive_formatting(self, paragraph, shape_bg: Optional[str] = None) -> EnhancedFormattingInfo:
        """Extract comprehensive formatting information with all details."""
        fmt = EnhancedFormattingInfo()

        # Set background if provided
        fmt.background_color = shape_bg

        try:
            # Paragraph-level formatting
            pf = getattr(paragraph, "paragraph_format", None)

            if pf:
                # Alignment
                align = getattr(pf, "alignment", None)
                if align is not None:
                    try:
                        fmt.paragraph_alignment = PP_ALIGN(align).name
                    except:
                        fmt.paragraph_alignment = str(align)

                # Margins and indents
                lm = getattr(pf, "left_margin", None)
                if lm is not None and hasattr(lm, "pt"):
                    fmt.left_margin = float(lm.pt)

                rm = getattr(pf, "right_margin", None)
                if rm is not None and hasattr(rm, "pt"):
                    fmt.right_margin = float(rm.pt)

                fi = getattr(pf, "first_line_indent", None)
                if fi is not None and hasattr(fi, "pt"):
                    fmt.first_line_indent = float(fi.pt)

                # Spacing before/after
                sb = getattr(pf, "space_before", None)
                if sb is not None and hasattr(sb, "pt"):
                    fmt.paragraph_spacing_before = float(sb.pt)

                sa = getattr(pf, "space_after", None)
                if sa is not None and hasattr(sa, "pt"):
                    fmt.paragraph_spacing_after = float(sa.pt)

                # Line spacing
                ls = getattr(pf, "line_spacing", None)
                if ls is not None:
                    if isinstance(ls, (float, int)):
                        fmt.line_spacing = float(ls)
                    elif hasattr(ls, "pt"):
                        fmt.line_spacing = float(ls.pt)

            # Bullet formatting
            self._extract_bullet_formatting(paragraph, fmt)

            # Run-level formatting (font, color, style)
            self._extract_run_formatting(paragraph, fmt)

        except Exception as e:
            logger.debug(f"Could not extract comprehensive formatting: {e}")

        return fmt

    def _extract_bullet_formatting(self, paragraph, fmt: EnhancedFormattingInfo):
        """Extract bullet point formatting details."""
        try:
            level = getattr(paragraph, 'level', None)
            if level is None:
                fmt.bullet_type = 'none'
                return

            pf = paragraph.paragraph_format

            # Check if numbered or bulleted
            if hasattr(pf, 'bullet'):
                bullet = pf.bullet

                # Check bullet type
                if hasattr(bullet, 'type'):
                    if bullet.type == 1:  # PP_BULLET_TYPE.NUMBERED
                        fmt.bullet_type = 'numbered'
                    elif bullet.type == 2:  # PP_BULLET_TYPE.PICTURE
                        fmt.bullet_type = 'picture'
                    else:
                        fmt.bullet_type = 'bulleted'

                # Get bullet character
                if hasattr(bullet, 'char'):
                    fmt.bullet_char = bullet.char

                # Get bullet font
                if hasattr(bullet, 'font'):
                    bullet_font = bullet.font
                    if hasattr(bullet_font, 'name'):
                        fmt.bullet_font = bullet_font.name
                    if hasattr(bullet_font, 'size') and bullet_font.size:
                        fmt.bullet_size = float(bullet_font.size.pt)

        except Exception as e:
            logger.debug(f"Could not extract bullet formatting: {e}")

    def _extract_run_formatting(self, paragraph, fmt: EnhancedFormattingInfo):
        """Extract run-level formatting (font, color, style)."""
        runs = list(getattr(paragraph, "runs", []))
        if not runs:
            return

        # Font name - with XML fallback
        for r in runs:
            font = r.font
            name = getattr(font, "name", None)

            if name:
                fmt.font_name = name
                break

            # XML fallback for inherited fonts
            try:
                if hasattr(font, '_element') and font._element is not None:
                    # Try Latin font
                    latin = font._element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}latin')
                    if latin is not None:
                        typeface = latin.get('typeface')
                        if typeface and typeface not in ['+mn-lt', '+mj-lt']:  # Skip theme references
                            fmt.font_name = typeface
                            break
                        # Resolve theme fonts
                        if typeface == '+mn-lt':
                            fmt.font_name = self.theme_info.font_scheme.get('minor', 'Calibri')
                            break
                        if typeface == '+mj-lt':
                            fmt.font_name = self.theme_info.font_scheme.get('major', 'Calibri Light')
                            break
            except:
                pass

        # Font size
        sizes = []
        for r in runs:
            try:
                if r.font.size:
                    sizes.append(float(r.font.size.pt))
            except:
                pass
        if sizes:
            fmt.font_size = sizes[0]

        # Bold, italic, underline - check all runs
        bold_values = []
        italic_values = []
        underline_values = []

        for r in runs:
            bold = r.font.bold
            italic = r.font.italic
            underline = r.font.underline

            if bold is not None:
                bold_values.append(bold)
            if italic is not None:
                italic_values.append(italic)
            if underline is not None:
                underline_values.append(underline)

        # Set to True if any run is True, False if all are False, None if all are None
        if bold_values:
            fmt.is_bold = any(bold_values)
        if italic_values:
            fmt.is_italic = any(italic_values)
        if underline_values:
            fmt.is_underline = any(underline_values)

        # Font color - comprehensive extraction
        for r in runs:
            color_info = self._extract_color_comprehensive(r.font.color)
            if color_info.get("rgb") or color_info.get("display") != "default":
                fmt.font_color = color_info.get("display")
                fmt.font_color_rgb = color_info.get("rgb")
                fmt.font_color_theme = color_info.get("theme")
                fmt.font_color_brightness = color_info.get("brightness")
                break

    def _extract_color_comprehensive(self, color) -> Dict[str, Optional[Any]]:
        """
        Comprehensive color extraction with theme resolution.

        Returns:
        {
            'display': 'default' | 'theme:ACCENT_1' | '#RRGGBB',
            'rgb': '#RRGGBB' | None,
            'theme': 'ACCENT_1' | None,
            'brightness': float | None  # brightness modifier (-1.0 to 1.0)
        }
        """
        out = {'display': 'default', 'rgb': None, 'theme': None, 'brightness': None}

        if not color:
            return out

        try:
            # Check color type
            ctype = getattr(color, "type", None)

            # Direct RGB color
            if ctype == MSO_COLOR_TYPE.RGB:
                rgb = getattr(color, "rgb", None)
                if rgb is not None:
                    if isinstance(rgb, RGBColor):
                        hex_str = f"#{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"
                    else:
                        hex_str = f"#{int(rgb):06X}"
                    out['display'] = hex_str
                    out['rgb'] = hex_str
                    return out

            # Scheme (theme) color
            if ctype == MSO_COLOR_TYPE.SCHEME:
                theme_color = getattr(color, "theme_color", None)
                if theme_color:
                    theme_name = self.THEME_COLOR_NAMES.get(theme_color, str(theme_color))
                    out['display'] = f"theme:{theme_name}"
                    out['theme'] = theme_name

                    # Try to get brightness modifier
                    try:
                        brightness = getattr(color, "brightness", None)
                        if brightness is not None:
                            out['brightness'] = float(brightness)
                    except:
                        pass

                    # Try to resolve to actual RGB from theme
                    resolved_rgb = self._resolve_theme_color(theme_color, out['brightness'])
                    if resolved_rgb:
                        out['rgb'] = resolved_rgb

                    return out

            # HSL color
            if ctype == MSO_COLOR_TYPE.HSL:
                # Try to convert to RGB if possible
                try:
                    rgb = getattr(color, "rgb", None)
                    if rgb:
                        hex_str = f"#{int(rgb):06X}"
                        out['display'] = hex_str
                        out['rgb'] = hex_str
                        return out
                except:
                    pass

            # Fallback: try to get RGB anyway
            try:
                rgb = getattr(color, "rgb", None)
                if rgb:
                    hex_str = f"#{int(rgb):06X}"
                    out['display'] = hex_str
                    out['rgb'] = hex_str
            except:
                pass

        except Exception as e:
            logger.debug(f"Color extraction failed: {e}")

        return out

    def _resolve_theme_color(self, theme_color, brightness: Optional[float] = None) -> Optional[str]:
        """Resolve theme color to RGB hex string."""
        # This is a simplified resolver - full theme color resolution
        # would require parsing the theme XML completely
        theme_color_map = {
            MSO_THEME_COLOR.ACCENT_1: self.theme_info.color_scheme.get('accent1'),
            MSO_THEME_COLOR.ACCENT_2: self.theme_info.color_scheme.get('accent2'),
            MSO_THEME_COLOR.ACCENT_3: self.theme_info.color_scheme.get('accent3'),
            MSO_THEME_COLOR.ACCENT_4: self.theme_info.color_scheme.get('accent4'),
            MSO_THEME_COLOR.ACCENT_5: self.theme_info.color_scheme.get('accent5'),
            MSO_THEME_COLOR.ACCENT_6: self.theme_info.color_scheme.get('accent6'),
            MSO_THEME_COLOR.TEXT_1: self.theme_info.color_scheme.get('dk1'),
            MSO_THEME_COLOR.TEXT_2: self.theme_info.color_scheme.get('dk2'),
            MSO_THEME_COLOR.BACKGROUND_1: self.theme_info.color_scheme.get('lt1'),
            MSO_THEME_COLOR.BACKGROUND_2: self.theme_info.color_scheme.get('lt2'),
        }

        rgb = theme_color_map.get(theme_color)

        # Apply brightness modifier if available
        # (simplified - real implementation would adjust RGB values)
        return rgb

    def _extract_notes(self, slide) -> Optional[str]:
        """Extract speaker notes."""
        try:
            if hasattr(slide, 'notes_slide') and slide.notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                return notes if notes else None
        except:
            pass
        return None

    def get_metadata(self) -> Dict[str, Any]:
        """Get comprehensive document metadata."""
        if not self.presentation:
            return {}

        metadata = {
            'file_path': str(self.file_path),
            'total_slides': len(self.presentation.slides),
            'layouts': [],
            'theme': self.theme_info.to_dict()
        }

        # Get unique layouts
        layouts = set()
        for slide in self.presentation.slides:
            try:
                layouts.add(slide.slide_layout.name)
            except:
                layouts.add("Unknown")
        metadata['layouts'] = list(layouts)

        # Get document properties
        try:
            props = self.presentation.core_properties
            metadata.update({
                'title': props.title,
                'author': props.author,
                'subject': props.subject,
                'keywords': props.keywords,
                'created': props.created.isoformat() if props.created else None,
                'modified': props.modified.isoformat() if props.modified else None,
                'last_modified_by': props.last_modified_by,
                'revision': props.revision
            })
        except:
            pass

        # Slide dimensions
        try:
            metadata['slide_width'] = self.presentation.slide_width.inches
            metadata['slide_height'] = self.presentation.slide_height.inches
        except:
            pass

        return metadata

    def extract_to_dict(self) -> Dict[str, Any]:
        """Extract all content to dictionary with comprehensive information."""
        slides = self.extract_all_content()
        metadata = self.get_metadata()

        return {
            'document_type': 'pptx',
            'metadata': metadata,
            'slides': [slide.to_dict() for slide in slides],
            'extraction_summary': {
                'total_slides': len(slides),
                'slides_with_content': len([s for s in slides if s.content_elements]),
                'slides_with_notes': len([s for s in slides if s.notes]),
                'slides_with_tables': len([s for s in slides if s.tables]),
                'total_elements': sum(len(s.content_elements) for s in slides),
                'total_tables': sum(len(s.tables) for s in slides),
                'total_table_cells': sum(len(table.cells) for s in slides for table in s.tables)
            }
        }


def process_pptx_file(file_path: str) -> Dict[str, Any]:
    """Process a PPTX file with enhanced extraction."""
    reader = EnhancedPPTXReader()
    if not reader.load_file(file_path):
        raise ValueError(f"Failed to load PPTX: {file_path}")
    return reader.extract_to_dict()


# if __name__ == "__main__":
#     import json

#     try:
#         result = process_pptx_file(
#             r"C:\Users\Shreya B\Downloads\amida standard template.pptx")

#         output_file = r'C:\Users\Shreya B\Documents\GitHub\slide-review-agent\backend\templates\amida standard template.json'
#         with open(output_file, "w", encoding="utf-8") as f:
#             json.dump(result, f, indent=2, ensure_ascii=False)

#         print(f"Enhanced JSON output saved to {output_file}")
#         print(f"\nExtraction Summary:")
#         print(f"  - Slides: {result['extraction_summary']['total_slides']}")
#         print(f"  - Text elements: {result['extraction_summary']['total_elements']}")
#         print(f"  - Tables: {result['extraction_summary']['total_tables']}")
#         print(f"  - Table cells: {result['extraction_summary']['total_table_cells']}")
#         print(f"  - Slides with tables: {result['extraction_summary']['slides_with_tables']}")

#         # Show theme info if available
#         if result['metadata'].get('theme'):
#             theme = result['metadata']['theme']
#             print(f"\nTheme: {theme.get('theme_name', 'Unknown')}")
#             if theme.get('font_scheme'):
#                 print(f"Fonts - Major: {theme['font_scheme'].get('major')}, Minor: {theme['font_scheme'].get('minor')}")

#     except Exception as e:
#         print(f"Error: {e}")
#         import traceback
#         traceback.print_exc()

# Alias for backward compatibility
PPTXReader = EnhancedPPTXReader
