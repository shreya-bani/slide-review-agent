"""
PowerPoint file reader for extracting text content with precise locators.
"""
import logging
from pathlib import Path
from typing import Dict, List, Optional, Any
from dataclasses import dataclass, asdict

from pptx import Presentation
from pptx.text.text import TextFrame
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_THEME_COLOR
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN


import warnings
warnings.filterwarnings("ignore", category=DeprecationWarning)

logger = logging.getLogger(__name__)


@dataclass
class TextLocator:
    """Precise location information for text elements."""
    slide_index: int
    element_type: str  # 'title', 'content', 'bullet', 'notes'
    element_index: int
    paragraph_index: Optional[int] = None
    bullet_level: Optional[int] = None
    
    def to_dict(self) -> Dict[str, Any]:
        return asdict(self)


@dataclass
class FormattingInfo:
    """Formatting information for text elements."""
    font_name: Optional[str] = None
    font_size: Optional[float] = None
    is_bold: Optional[bool] = None
    is_italic: Optional[bool] = None
    # Font color
    font_color: Optional[str] = None
    font_color_rgb: Optional[str] = None  # Resolved RGB hex if available
    # Paragraph spacing
    paragraph_spacing_before: Optional[float] = None  # points
    paragraph_spacing_after: Optional[float] = None   # points
    line_spacing: Optional[float] = None              # multiple (float) or exact (points)
    # Extra paragraph layout fields
    paragraph_alignment: Optional[str] = None         # e.g. 'LEFT', 'CENTER', 'RIGHT', 'JUSTIFY'
    left_margin: Optional[float] = None               # points
    first_line_indent: Optional[float] = None         # points
    def to_dict(self) -> Dict[str, Any]:
        return asdict(self)



@dataclass
class ExtractedText:
    """Text content with location and formatting information."""
    text: str
    locator: TextLocator
    formatting: FormattingInfo
    
    def to_dict(self) -> Dict[str, Any]:
        return {
            'text': self.text,
            'locator': self.locator.to_dict(),
            'formatting': self.formatting.to_dict()
        }


@dataclass
class SlideContent:
    """Complete content from a single slide."""
    slide_index: int
    title: Optional[str] = None
    content_elements: List[ExtractedText] = None
    notes: Optional[str] = None
    layout_name: Optional[str] = None
    
    def __post_init__(self):
        if self.content_elements is None:
            self.content_elements = []
    
    def to_dict(self) -> Dict[str, Any]:
        return {
            'slide_index': self.slide_index,
            'title': self.title,
            'content_elements': [elem.to_dict() for elem in self.content_elements],
            'notes': self.notes,
            'layout_name': self.layout_name,
            'element_count': len(self.content_elements)
        }


class PPTXReader:
    """Extract text content from PowerPoint files."""
    
    # Theme color name mapping
    THEME_COLOR_NAMES = {
        MSO_THEME_COLOR.ACCENT_1: "ACCENT_1",
        MSO_THEME_COLOR.ACCENT_2: "ACCENT_2",
        MSO_THEME_COLOR.ACCENT_3: "ACCENT_3",
        MSO_THEME_COLOR.ACCENT_4: "ACCENT_4",
        MSO_THEME_COLOR.ACCENT_5: "ACCENT_5",
        MSO_THEME_COLOR.ACCENT_6: "ACCENT_6",
        MSO_THEME_COLOR.BACKGROUND_1: "BACKGROUND_1",
        MSO_THEME_COLOR.BACKGROUND_2: "BACKGROUND_2",
        MSO_THEME_COLOR.DARK_1: "DARK_1",
        MSO_THEME_COLOR.DARK_2: "DARK_2",
        MSO_THEME_COLOR.FOLLOWED_HYPERLINK: "FOLLOWED_HYPERLINK",
        MSO_THEME_COLOR.HYPERLINK: "HYPERLINK",
        MSO_THEME_COLOR.LIGHT_1: "LIGHT_1",
        MSO_THEME_COLOR.LIGHT_2: "LIGHT_2",
        MSO_THEME_COLOR.TEXT_1: "TEXT_1",
        MSO_THEME_COLOR.TEXT_2: "TEXT_2",
    }
    
    def __init__(self):
        self.presentation = None
        self.file_path = None
        self.theme_colors = {}
    
    def load_file(self, file_path: str) -> bool:
        """Load PowerPoint file."""
        try:
            self.file_path = Path(file_path)
            self.presentation = Presentation(file_path)
            self._extract_theme_colors()
            logger.info(f"Loaded PPTX: {file_path} ({len(self.presentation.slides)} slides)")
            return True
        except Exception as e:
            logger.error(f"Failed to load PPTX {file_path}: {e}")
            return False
    
    def _extract_theme_colors(self):
        """Extract theme colors from the presentation."""
        try:
            if hasattr(self.presentation, 'part') and hasattr(self.presentation.part, 'theme'):
                theme = self.presentation.part.theme
                if hasattr(theme, 'color_scheme'):
                    # Store reference to color scheme for potential future use
                    self.theme_colors['_scheme'] = theme.color_scheme
        except Exception as e:
            logger.debug(f"Could not extract theme colors: {e}")
    
    def extract_all_content(self) -> List[SlideContent]:
        """Extract content from all slides (1-based indexing)."""
        if not self.presentation:
            raise ValueError("No presentation loaded")
        
        slides = []
        for idx, slide in enumerate(self.presentation.slides, start=1):
            slide_content = self._extract_slide_content(idx, slide)
            slides.append(slide_content)
        
        return slides
    
    def _extract_slide_content(self, slide_idx: int, slide) -> SlideContent:
        """Extract content from a single slide."""
        slide_content = SlideContent(slide_index=slide_idx)
        
        # Get layout name
        try:
            slide_content.layout_name = slide.slide_layout.name
        except:
            slide_content.layout_name = "Unknown"
        
        # Extract title
        slide_content.title = self._extract_title(slide)
        
        # Extract content from shapes
        slide_content.content_elements = self._extract_shapes_content(slide_idx, slide)
        
        # Extract notes
        slide_content.notes = self._extract_notes(slide)
        
        return slide_content
    
    def _extract_title(self, slide) -> Optional[str]:
        """Extract slide title."""
        try:
            if hasattr(slide, 'shapes') and slide.shapes.title:
                title = slide.shapes.title.text.strip()
                return title if title else None
        except:
            pass
        return None
    
    def _extract_shapes_content(self, slide_idx: int, slide) -> List[ExtractedText]:
        """Extract text from all shapes in slide."""
        elements = []
        content_idx = 0
        
        for shape in slide.shapes:
            if not hasattr(shape, 'text_frame') or not shape.text_frame:
                continue
            
            # Skip title (already extracted)
            if shape == slide.shapes.title:
                continue
            
            shape_elements = self._extract_text_frame(slide_idx, shape.text_frame, content_idx)
            elements.extend(shape_elements)
            content_idx += 1
        
        return elements
    
    def _extract_text_frame(self, slide_idx: int, text_frame: TextFrame, 
                           shape_idx: int) -> List[ExtractedText]:
        """Extract text from a text frame."""
        elements = []
        
        for para_idx, paragraph in enumerate(text_frame.paragraphs):
            text = paragraph.text.strip()
            if not text:
                continue
            
            # Get bullet level
            bullet_level = getattr(paragraph, 'level', None)
            
            # Get formatting info
            formatting = self._extract_formatting(paragraph)
            
            # Determine element type
            element_type = 'bullet' if bullet_level is not None else 'content'
            
            # Create locator
            locator = TextLocator(
                slide_index=slide_idx,
                element_type=element_type,
                element_index=shape_idx,
                paragraph_index=para_idx,
                bullet_level=bullet_level
            )
            
            # Create extracted text
            element = ExtractedText(
                text=text,
                locator=locator,
                formatting=formatting
            )
            
            elements.append(element)
        
        return elements
    
   
            
    def _extract_formatting(self, paragraph) -> FormattingInfo:
        """Extract formatting information from paragraph with enhanced fallback."""
        fmt = FormattingInfo()

        try:
            # Paragraph-level spacing
            pf = getattr(paragraph, "paragraph_format", None)

            align = getattr(pf, "alignment", None)
            if align is not None:
                # store the enum name, e.g. 'LEFT', 'CENTER'
                try:
                    fmt.paragraph_alignment = PP_ALIGN(align).name  # add a field to FormattingInfo if you want this
                except Exception:
                    fmt.paragraph_alignment = str(align)

            lm = getattr(pf, "left_margin", None)
            if lm is not None and hasattr(lm, "pt"):
                fmt.left_margin = float(lm.pt)  # add fields in FormattingInfo if you want these

            fi = getattr(pf, "first_line_indent", None)
            if fi is not None and hasattr(fi, "pt"):
                fmt.first_line_indent = float(fi.pt)

            # spacing before/after (points)
            if pf is not None:
                sb = getattr(pf, "space_before", None)
                if sb is not None and hasattr(sb, "pt"):
                    fmt.paragraph_spacing_before = float(sb.pt)

                sa = getattr(pf, "space_after", None)
                if sa is not None and hasattr(sa, "pt"):
                    fmt.paragraph_spacing_after = float(sa.pt)

                # line spacing: float (multiplier) OR Length (exact pts)
                ls = getattr(pf, "line_spacing", None)
                if ls is not None:
                    # float → store multiplier as-is
                    if isinstance(ls, float) or isinstance(ls, int):
                        fmt.line_spacing = float(ls)
                    else:
                        # likely a Length – keep points
                        pt = getattr(ls, "pt", None)
                        if pt is not None:
                            fmt.line_spacing = float(pt)

            # Run-level font info
            runs = list(getattr(paragraph, "runs", []))
            if runs:
                # font family: take the first non-None across runs
                # Try multiple approaches to get font name
                for r in runs:
                    font = r.font
                    # First try direct name access
                    name = getattr(font, "name", None)
                    if name:
                        fmt.font_name = name
                        break
                    # If name is None, check if there's an element with font family
                    try:
                        if hasattr(font, '_element') and font._element is not None:
                            latin = font._element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}latin')
                            if latin is not None and latin.get('typeface'):
                                fmt.font_name = latin.get('typeface')
                                break
                    except Exception:
                        pass

                # font size: collect all defined sizes (points)
                sizes = []
                for r in runs:
                    f = r.font
                    try:
                        if getattr(f, "size", None):
                            sizes.append(float(f.size.pt))
                    except Exception:
                        pass

                # choose your policy: first non-None, or max, etc.
                # here we pick the first defined size for "effective" size
                if sizes:
                    fmt.font_size = sizes[0]

                # bold/italic: if any run is explicitly True, mark True;
                # if all None/False, stays None/False
                is_bold = None
                is_italic = None
                for r in runs:
                    b = r.font.bold
                    i = r.font.italic
                    if b is True:
                        is_bold = True
                    if i is True:
                        is_italic = True
                fmt.is_bold = is_bold
                fmt.is_italic = is_italic

                # color: resolve from the first run that yields something meaningful
                for r in runs:
                    color_info = self._extract_color_enhanced(r.font, paragraph)
                    if color_info.get("rgb") or color_info.get("display") != "default":
                        fmt.font_color = color_info.get("display")
                        fmt.font_color_rgb = color_info.get("rgb")
                        break

        except Exception as e:
            logger.debug(f"Could not extract formatting: {e}")

        return fmt

    def _extract_color_enhanced(self, font, paragraph) -> Dict[str, Optional[str]]:
        """
        Robust color extraction.
        Returns:
        {
            'display': 'default' | 'theme:ACCENT_1' | '#RRGGBB',
            'rgb': '#RRGGBB' | None
        }
        """
        out = {'display': 'default', 'rgb': None}

        color = getattr(font, "color", None)
        if not color:
            return out

        try:
            # Direct RGB?
            # color.rgb is an RGBColor (int-like); convert to hex.
            rgb = getattr(color, "rgb", None)
            if rgb is not None:
                hex_str = f"#{int(rgb):06X}"
                out['display'] = hex_str
                out['rgb'] = hex_str
                return out

            # Scheme (theme) color?
            ctype = getattr(color, "type", None)
            if ctype == MSO_COLOR_TYPE.SCHEME:
                theme_color = getattr(color, "theme_color", None)
                theme_name = self.THEME_COLOR_NAMES.get(theme_color, str(theme_color))
                out['display'] = f"theme:{theme_name}"
                # Some themes also expose brightness/tint; python-pptx doesn't resolve final RGB reliably.
                # You could attempt manual theme resolution here if you need exact RGB.
                return out

            # Other types: attempt hex if possible
            hex_str = None
            rgb = getattr(color, "rgb", None)
            if rgb is not None:
                hex_str = f"#{int(rgb):06X}"
            if hex_str:
                out['display'] = hex_str
                out['rgb'] = hex_str
            else:
                out['display'] = f"color_type:{ctype}"
            return out

        except Exception as e:
            logger.debug(f"Color extraction failed: {e}")
            return out
        

    def _extract_notes(self, slide) -> Optional[str]:
        """Extract speaker notes."""
        try:
            if hasattr(slide, 'notes_slide') and slide.notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                return notes if notes else None
        except:
            pass
        return None
    
    def get_metadata(self) -> Dict[str, Any]:
        """Get document metadata."""
        if not self.presentation:
            return {}
        
        metadata = {
            'file_path': str(self.file_path),
            'total_slides': len(self.presentation.slides),
            'layouts': []
        }
        
        # Get unique layouts
        layouts = set()
        for slide in self.presentation.slides:
            try:
                layouts.add(slide.slide_layout.name)
            except:
                layouts.add("Unknown")
        metadata['layouts'] = list(layouts)
        
        # Get document properties
        try:
            props = self.presentation.core_properties
            metadata.update({
                'title': props.title,
                'author': props.author,
                'created': props.created.isoformat() if props.created else None,
                'modified': props.modified.isoformat() if props.modified else None
            })
        except:
            pass
        
        return metadata
    
    def extract_to_dict(self) -> Dict[str, Any]:
        """Extract all content to dictionary."""
        slides = self.extract_all_content()
        metadata = self.get_metadata()
        
        return {
            'document_type': 'pptx',
            'metadata': metadata,
            'slides': [slide.to_dict() for slide in slides],
            'extraction_summary': {
                'total_slides': len(slides),
                'slides_with_content': len([s for s in slides if s.content_elements]),
                'slides_with_notes': len([s for s in slides if s.notes]),
                'total_elements': sum(len(s.content_elements) for s in slides)
            }
        }


def process_pptx_file(file_path: str) -> Dict[str, Any]:
    """Process a PPTX file and return content."""
    reader = PPTXReader()
    if not reader.load_file(file_path):
        raise ValueError(f"Failed to load PPTX: {file_path}")
    return reader.extract_to_dict()


if __name__ == "__main__":
    import json

    try:
        result = process_pptx_file(
            r"C:\Users\Shreya B\Documents\GitHub\slide-review-agent\data\uploads\001_Amida_Agentic_AI_solution_Strategic_Plan_Draft_Aug.pptx"
        )

        output_file = r'C:\Users\Shreya B\Documents\GitHub\slide-review-agent\data\logs\try.json'
        with open(output_file, "w", encoding="utf-8") as f:
            json.dump(result, f, indent=2, ensure_ascii=False)

        print(f"JSON output saved to {output_file}")

    except Exception as e:
        print(f"Error: {e}")
        import traceback
        traceback.print_exc()