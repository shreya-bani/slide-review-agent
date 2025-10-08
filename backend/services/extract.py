# --- compatibility shim for python-pptx / collections on some envs ---
import collections  # noqa: F401
import collections.abc  # noqa: F401

import argparse
import csv
import glob
import os
from pathlib import Path
from typing import List, Iterable

# Optional deps: only needed when processing these types
try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except Exception:
    Presentation = None
    MSO_SHAPE_TYPE = None

try:
    import fitz  # PyMuPDF
except Exception:
    fitz = None


CSV_HEADER = ["File", "Page", "Text", "Notes", "Images"]


def safe_mkdir(path: Path) -> None:
    path.mkdir(parents=True, exist_ok=True)


def write_csv_header(csv_path: Path) -> None:
    if not csv_path.exists():
        with csv_path.open("w", encoding="utf-8", newline="") as f:
            csv.writer(f).writerow(CSV_HEADER)


def append_row(csv_path: Path, row: List):
    with csv_path.open("a", encoding="utf-8", newline="") as f:
        csv.writer(f).writerow(row)


# ----------------------- PPTX processing -----------------------
def _pptx_collect_text(slide) -> str:
    text_chunks: List[str] = []
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False) and shape.text and shape.text.strip():
            text_chunks.append(shape.text)
    return os.linesep.join(text_chunks)


def _pptx_save_images(slide, base_stem: str, images_dir: Path) -> List[str]:
    """
    Walk shapes (including groups) and save images as:
      images/<base_stem> p<page> <idx>.<ext>
    """
    saved: List[str] = []

    def drill(shape, img_index_start: int) -> int:
        idx = img_index_start
        if MSO_SHAPE_TYPE and shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for s in shape.shapes:
                idx = drill(s, idx)
        elif MSO_SHAPE_TYPE and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            # NOTE: page index is handled by caller in filename stem; here we just return name
            image = shape.image
            ext = image.ext or "png"
            filename = f"{base_stem} {idx}.{ext}"
            out_path = images_dir / filename
            with out_path.open("wb") as f:
                f.write(image.blob)
            saved.append(filename)
            idx += 1
        return idx

    # Start at 1 per slide
    img_idx = 1
    for shape in slide.shapes:
        img_idx = drill(shape, img_idx)
    return saved


def process_pptx_file(pptx_path: Path, images_dir: Path, csv_path: Path):
    if Presentation is None:
        raise RuntimeError("python-pptx is not installed. Run: pip install python-pptx")

    pres = Presentation(str(pptx_path))
    base = pptx_path.stem

    for page_idx, slide in enumerate(pres.slides, start=1):
        text = _pptx_collect_text(slide)

        # safe notes
        notes = ""
        if getattr(slide, "notes_slide", None) and getattr(slide.notes_slide, "notes_text_frame", None):
            notes = slide.notes_slide.notes_text_frame.text or ""

        # images
        # build a per-slide base stem so filenames include page number
        slide_stem = f"{base} p{page_idx}"
        images = _pptx_save_images(slide, slide_stem, images_dir)

        append_row(csv_path, [str(pptx_path), page_idx, text, notes, images])


# ----------------------- PDF processing -----------------------
def process_pdf_file(pdf_path: Path, images_dir: Path, csv_path: Path):
    if fitz is None:
        raise RuntimeError("PyMuPDF is not installed. Run: pip install pymupdf")

    doc = fitz.open(str(pdf_path))
    base = pdf_path.stem

    for page_idx, page in enumerate(doc, start=1):
        text = page.get_text("text") or ""
        images: List[str] = []

        for img_idx, img in enumerate(page.get_images(full=True), start=1):
            xref = img[0]
            meta = doc.extract_image(xref)
            ext = meta.get("ext", "png")
            data = meta["image"]

            filename = f"{base} p{page_idx} {img_idx}.{ext}"
            (images_dir / filename).write_bytes(data)
            images.append(filename)

        # PDFs have no speaker notes
        append_row(csv_path, [str(pdf_path), page_idx, text, "", images])


# ----------------------- Router & CLI -----------------------
def iter_targets(input_path: Path) -> Iterable[Path]:
    if input_path.is_file():
        yield input_path
    else:
        # process all .pptx and .pdf in folder (non-recursive)
        for ext in ("*.pptx", "*.pdf"):
            for p in input_path.glob(ext):
                yield p


def main():
    parser = argparse.ArgumentParser(
        description="Extract text and images from PPTX/PDF into text.csv and images/."
    )
    parser.add_argument("path", help="Path to a .pptx, .pdf, or a directory containing them.")
    parser.add_argument("--csv", default="text.csv", help="CSV output path (default: text.csv)")
    parser.add_argument("--images-dir", default="images", help="Directory to save images (default: images)")
    args = parser.parse_args()

    input_path = Path(args.path).expanduser().resolve()
    csv_path = Path(args.csv).expanduser().resolve()
    images_dir = Path(args.images_dir).expanduser().resolve()

    if not input_path.exists():
        raise FileNotFoundError(f"Input path not found: {input_path}")

    safe_mkdir(images_dir)
    write_csv_header(csv_path)

    processed = 0
    for target in iter_targets(input_path):
        suffix = target.suffix.lower()
        if suffix == ".pptx":
            process_pptx_file(target, images_dir, csv_path)
            processed += 1
        elif suffix == ".pdf":
            process_pdf_file(target, images_dir, csv_path)
            processed += 1
        else:
            # skip unknown types silently
            continue

    print(f"Done. Processed {processed} file(s).")
    print(f"CSV: {csv_path}")
    print(f"Images dir: {images_dir}")


if __name__ == "__main__":
    main()
