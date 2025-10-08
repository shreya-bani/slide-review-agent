#!/usr/bin/env python3
r"""
pptx_extract_to_json_md.py

Simplified usage:
  python pptx_extract_to_json_md.py "C:/path/to/slides.pptx"

Outputs:
  JSON  -> C:\Users\Shreya B\Documents\GitHub\slide-review-agent\data\logs\<filename>.json
  Images-> ...\data\logs\assets\<deckname>\

Requires:
  pip install python-pptx
"""

from __future__ import annotations
import json, sys
from dataclasses import asdict, dataclass, field
from pathlib import Path
from typing import Any, Dict, List, Optional, Union
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# ---------------- Data Models ----------------

@dataclass
class FontInfo:
    name: Optional[str] = None
    size_pt: Optional[float] = None
    bold: Optional[bool] = None
    italic: Optional[bool] = None
    underline: Optional[bool] = None
    # Color details (robust to theme colors)
    color_rgb: Optional[str] = None        # '#RRGGBB' when concrete RGB is present
    color_theme: Optional[str] = None      # e.g. 'ACCENT_1', 'DARK_1', etc.
    color_brightness: Optional[float] = None

@dataclass
class TextRun:
    text: str
    font: FontInfo = field(default_factory=FontInfo)

@dataclass
class Paragraph:
    runs: List[TextRun] = field(default_factory=list)
    level: int = 0  # bullet level if available

@dataclass
class TextBoxElement:
    element_index: int
    shape_id: int
    placeholder_type: Optional[str]
    kind: str  # "textbox" | "title" | "subtitle" | "footer" | ...
    bbox: Dict[str, int]  # approximate pixels
    paragraphs: List[Paragraph] = field(default_factory=list)

@dataclass
class ImageElement:
    element_index: int
    shape_id: int
    kind: str  # "image"
    bbox: Dict[str, int]
    filename: str  # saved file path
    image_ext: str
    image_width_px: Optional[int] = None
    image_height_px: Optional[int] = None

@dataclass
class TableCell:
    row: int
    col: int
    text: str

@dataclass
class TableElement:
    element_index: int
    shape_id: int
    kind: str  # "table"
    bbox: Dict[str, int]
    rows: int
    cols: int
    cells: List[TableCell] = field(default_factory=list)

@dataclass
class SlideData:
    slide_index: int  # 1-based for humans
    title: Optional[str]
    notes: Optional[str]
    elements: List[Union[TextBoxElement, ImageElement, TableElement]] = field(default_factory=list)

@dataclass
class DeckData:
    source: str
    slides: List[SlideData] = field(default_factory=list)
    meta: Dict[str, Any] = field(default_factory=dict)  # deck-level metadata (e.g., slide size)

# ---------------- Helpers ----------------

EMU_PER_INCH = 914400
DPI = 96  # approximation for mapping EMU bbox to pixels

def emu_to_px(emu: int) -> int:
    return int(round((emu / EMU_PER_INCH) * DPI))

def _shape_bbox(shape) -> Dict[str, int]:
    # Centralized bbox conversion; handles odd shapes safely.
    try:
        l, t, w, h = shape.left, shape.top, shape.width, shape.height
        return {"left": emu_to_px(l), "top": emu_to_px(t), "width": emu_to_px(w), "height": emu_to_px(h)}
    except Exception:
        return {"left": 0, "top": 0, "width": 0, "height": 0}

def _placeholder_name(shape) -> Optional[str]:
    # Normalize placeholder type to a clean token like "TITLE", "SUBTITLE", "BODY"
    try:
        ph = shape.placeholder_format
        if ph and ph.type is not None:
            raw = str(ph.type)  # 'MSO_PLACEHOLDER.TITLE' or 'TITLE'
            return raw.split(".")[-1].upper()
    except Exception:
        pass
    return None

def _infer_kind(shape, placeholder: Optional[str]) -> str:
    if placeholder:
        low = placeholder.lower()
        if "title" in low:
            return "title"
        if "subtitle" in low:
            return "subtitle"
        if "footer" in low:
            return "footer"
        if "body" in low:
            return "textbox"
    st = shape.shape_type
    if st == MSO_SHAPE_TYPE.PICTURE:
        return "image"
    if st == MSO_SHAPE_TYPE.TABLE:
        return "table"
    return "textbox"

def _font_color_triplet(color):
    """
    Safely extract color as (rgb_hex, theme_name, brightness).
    - rgb_hex: '#RRGGBB' if concrete RGB, else None
    - theme_name: 'ACCENT_1'...'ACCENT_6', 'DARK_1', etc., or None
    - brightness: float or None
    """
    rgb_hex = theme_name = brightness = None
    if color is None:
        return rgb_hex, theme_name, brightness
    try:
        rgb = getattr(color, "rgb", None)
        if rgb is not None:
            rgb_hex = f"#{int(rgb):06X}"
    except Exception:
        pass
    try:
        tc = getattr(color, "theme_color", None)
        if tc is not None:
            theme_name = str(tc)
    except Exception:
        pass
    try:
        brightness = getattr(color, "brightness", None)
    except Exception:
        pass
    return rgb_hex, theme_name, brightness

def _extract_textbox(shape, element_index: int) -> Optional[TextBoxElement]:
    if not getattr(shape, "has_text_frame", False):
        return None
    paragraphs: List[Paragraph] = []
    try:
        tframe = shape.text_frame
    except Exception:
        tframe = None
    if not tframe:
        return None

    for p in tframe.paragraphs:
        runs: List[TextRun] = []
        for r in p.runs:
            f = r.font
            rgb_hex, theme, bright = _font_color_triplet(getattr(f, "color", None))
            size_pt: Optional[float] = None
            try:
                size_pt = f.size.pt if f.size else None  # can be None (theme-inherited)
            except Exception:
                size_pt = None
            runs.append(
                TextRun(
                    r.text or "",
                    FontInfo(
                        name=getattr(f, "name", None),
                        size_pt=size_pt,
                        bold=(f.bold if f.bold is not None else None),
                        italic=(f.italic if f.italic is not None else None),
                        underline=(f.underline if f.underline is not None else None),
                        color_rgb=rgb_hex,
                        color_theme=theme,
                        color_brightness=bright,
                    ),
                )
            )
        paragraphs.append(Paragraph(runs=runs, level=getattr(p, "level", 0) or 0))

    ph_name = _placeholder_name(shape)
    element = TextBoxElement(
        element_index=element_index,
        shape_id=int(getattr(shape, "shape_id", element_index)),
        placeholder_type=ph_name,
        kind=_infer_kind(shape, ph_name),
        bbox=_shape_bbox(shape),
        paragraphs=paragraphs,
    )
    return element

def _save_picture(shape, slide_idx: int, assets_dir: Path, element_index: int) -> Optional[ImageElement]:
    try:
        image = shape.image
    except Exception:
        return None
    ext = image.ext or "png"
    out_path = assets_dir / f"{slide_idx:02d}_{element_index:02d}.{ext}"
    out_path.parent.mkdir(parents=True, exist_ok=True)
    with open(out_path, "wb") as f:
        f.write(image.blob)
    return ImageElement(
        element_index=element_index,
        shape_id=int(getattr(shape, "shape_id", element_index)),
        kind="image",
        bbox=_shape_bbox(shape),
        filename=str(out_path),
        image_ext=ext,
        image_width_px=getattr(image, "width", None),
        image_height_px=getattr(image, "height", None),
    )

def _extract_table(shape, element_index: int) -> Optional[TableElement]:
    try:
        tbl = shape.table
    except Exception:
        return None
    rows, cols = len(tbl.rows), len(tbl.columns)
    cells: List[TableCell] = []
    for r in range(rows):
        for c in range(cols):
            try:
                text = tbl.cell(r, c).text or ""
            except Exception:
                text = ""
            cells.append(TableCell(r, c, text))
    return TableElement(
        element_index=element_index,
        shape_id=int(getattr(shape, "shape_id", element_index)),
        kind="table",
        bbox=_shape_bbox(shape),
        rows=rows,
        cols=cols,
        cells=cells,
    )

def _slide_title(slide) -> Optional[str]:
    try:
        return slide.shapes.title.text if slide.shapes.title else None
    except Exception:
        return None

def _slide_notes(slide) -> Optional[str]:
    try:
        if slide.has_notes_slide and slide.notes_slide and slide.notes_slide.notes_text_frame:
            return slide.notes_slide.notes_text_frame.text
    except Exception:
        pass
    return None

# ---------------- Extraction ----------------

def extract_deck(pptx_path: Path, assets_dir: Path) -> DeckData:
    prs = Presentation(pptx_path.as_posix())
    deck = DeckData(str(pptx_path.resolve()))
    deck.meta = {
        "slide_width_px": emu_to_px(prs.slide_width),
        "slide_height_px": emu_to_px(prs.slide_height),
    }

    assets_dir.mkdir(parents=True, exist_ok=True)

    for i, slide in enumerate(prs.slides, 1):
        sd = SlideData(
            slide_index=i,
            title=_slide_title(slide),
            notes=_slide_notes(slide),
            elements=[],
        )
        idx = 0
        for shape in slide.shapes:
            idx += 1
            st = shape.shape_type
            if st == MSO_SHAPE_TYPE.PICTURE:
                img = _save_picture(shape, i, assets_dir, idx)
                if img:
                    sd.elements.append(img)
                continue
            if st == MSO_SHAPE_TYPE.TABLE:
                tbl = _extract_table(shape, idx)
                if tbl:
                    sd.elements.append(tbl)
                continue
            if getattr(shape, "has_text_frame", False) or st in (MSO_SHAPE_TYPE.TEXT_BOX, MSO_SHAPE_TYPE.PLACEHOLDER):
                tb = _extract_textbox(shape, idx)
                if tb:
                    sd.elements.append(tb)
                continue
            # Fallback: try to capture text anyway if available
            try:
                if getattr(shape, "has_text_frame", False):
                    tb = _extract_textbox(shape, idx)
                    if tb:
                        sd.elements.append(tb)
            except Exception:
                pass

        deck.slides.append(sd)
    return deck

# ---------------- IO ----------------

def save_json(deck: DeckData, out_path: Path) -> None:
    out_path.parent.mkdir(parents=True, exist_ok=True)
    with open(out_path, "w", encoding="utf-8") as f:
        json.dump(asdict(deck), f, indent=2, ensure_ascii=False)

# ---------------- Main ----------------

def main():
    # Single required arg: PPTX path
    if len(sys.argv) != 2 or not sys.argv[1].lower().endswith(".pptx"):
        print("Usage: python pptx_extract_to_json_md.py <path-to-pptx>")
        sys.exit(1)

    # Normalize path; forward slashes are OK in Windows too.
    pptx_arg = sys.argv[1]
    pptx_path = Path(pptx_arg)

    # Output locations (fixed)
    out_dir = Path(r"C:\Users\Shreya B\Documents\GitHub\slide-review-agent\data\logs")
    out_path = out_dir / (pptx_path.stem + ".json")
    assets_dir = out_dir / "assets" / pptx_path.stem

    deck = extract_deck(pptx_path, assets_dir)
    save_json(deck, out_path)
    print(f"✓ Extracted {pptx_path.name} → {out_path}")

if __name__ == "__main__":
    main()
